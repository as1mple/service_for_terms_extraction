import nltk
from pyate.term_extraction_pipeline import TermExtractionPipeline
from nltk.tokenize import word_tokenize
from sklearn.feature_extraction.text import TfidfVectorizer

from nltk.stem import WordNetLemmatizer
from nltk.corpus import stopwords
import pandas as pd
import spacy
import re
import numpy as np
import zipfile
import os

from openpyxl import load_workbook
from pptx import Presentation

import textract
from bs4 import BeautifulSoup
from docx2txt import docx2txt
import extract_msg
from odf import text, teletype
from odf.opendocument import load
from email import policy
from email.parser import BytesParser
import glob
import PyPDF2


class Extract:
    def __init__(self, path: str):
        self.UNIT_TO_MULTIPLIER = {

            '.csv': self.get_text_with_txt,
            '.doc': self.get_text_with_docx,
            '.docx': self.get_text_with_docx,
            '.eml': self.get_text_with_eml,
            '.html': self.get_text_with_html,
            '.msg': self.get_text_with_msg,
            '.pdf': self.get_text_with_pdf,
            '.pptx': self.get_text_with_pptx,
            '.txt': self.get_text_with_txt,
            '.xlsx': self.get_text_with_xml,
            '.xls': self.get_text_with_xml,
            '.idml': None,
            '.mht': None,
            '.odt': self.get_text_with_odt,
            '.odp': None,
            '.ods': None,
            '.xml': None,
            '.xps': None,
            '.zip': None
        }
        self.path = path

    def extract(self):
        return self.switch_case()

    def get_format(self):
        filename, file_extension = os.path.splitext(self.path)
        return file_extension

    def switch_case(self) -> str:
        try:

            command = self.UNIT_TO_MULTIPLIER[self.get_format()]
            return command() if command is not None else "This file format is not supported"

        except Exception as e:
            print(e)
            return None

    def get_text_with_eml(self) -> str:

        file_list = glob.glob('*.eml')  # returns list of files
        with open(file_list[2], 'rb') as fp:  # select a specific email file from the list
            msg = BytesParser(policy=policy.default).parse(fp)
        return msg.get_body(preferencelist=('plain')).get_content()

    def get_text_with_xml(self) -> str:
        workbook = load_workbook(self.path)
        first_sheet = workbook.get_sheet_names()[0]
        worksheet = workbook.get_sheet_by_name(first_sheet)
        text = ""
        for row in worksheet.iter_rows():
            text += str(row)

        # check out the last row
        for cell in row:
            text += str(cell)
        return text

    def get_text_with_docx(self) -> str:
        return docx2txt.process(self.path)

    def get_text_with_pptx(self):
        result = ""
        for eachfile in glob.glob(self.path):
            prs = Presentation(eachfile)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        result += shape.text
        return result

    def get_text_with_txt(self) -> str:
        return textract.process(self.path)

    def get_text_with_html(self) -> str:
        with open(self.path) as fp:
            soup = BeautifulSoup(fp, features="html.parser")

        # kill all script and style elements
        for script in soup(["script", "style"]):
            script.extract()  # rip it out

        # get text
        text = soup.get_text()

        # break into lines and remove leading and trailing space on each
        lines = (line.strip() for line in text.splitlines())
        # break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # drop blank lines
        text = '\n'.join(chunk for chunk in chunks if chunk)

        return text

    def get_text_with_msg(self) -> str:
        return extract_msg.openMsg(self.path)

    def get_text_with_odt(self) -> str:
        textdoc = load(self.path)
        allparas = textdoc.getElementsByType(text.P)
        return teletype.extractText(allparas)

    def get_text_with_pdf(self) -> str:
        # creating a pdf file object
        pdfFileObj = open('text/test.pdf', 'rb')

        # creating a pdf reader object
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
        text = ""
        for i in range(int(pdfReader.numPages)):
            # creating a page object
            pageObj = pdfReader.getPage(i)

            # extracting text from page
            text += pageObj.extractText()

        # closing the pdf file object
        pdfFileObj.close()
        return text


class TermsExtraction:

    def __init__(self, directory, nlp: spacy):
        self.directory = directory
        self.nlp = nlp
        self.nlp.add_pipe(TermExtractionPipeline())
        self.result_extract_terms = {}  # рузультат извлечение терминов из текста
        self.cleaning_dict_terms = {}  # результат удаление корелирующих фраз
        self.tf_ifd_score = {}
        self.corpus_dict = {}

    def get_terms_dict_with_text(self, text) -> dict:
        doc = self.nlp(text)
        return dict(doc._.combo_basic.sort_values(ascending=False))

    def get_clean_terms(self):
        return self.cleaning_dict_terms

    def get_extract_terms(self):
        return self.result_extract_terms

    @staticmethod
    def tf_idf(corpus, res):
        vectorizer = TfidfVectorizer(ngram_range=(1, 5))
        X2 = vectorizer.fit_transform(list(corpus.values()))
        dictir = {list(corpus.keys())[i]: X2[i].T.todense() for i in range(X2.shape[0])}
        df = pd.DataFrame({'terms': vectorizer.get_feature_names()})
        df['terms_score'] = ['NaN'] * df.shape[0]
        for k, v in dictir.items():
            df[k] = v

        for key, value in res.items():
            if len(df[df.terms == key]) > 0:
                df.loc[df['terms'] == key, ['terms_score']] = value['weight']
        return df.drop(np.where(df['terms_score'] == 'NaN')[0])

    @staticmethod
    def update_score(global_dict: dict, key: str, value, index_doc):
        last_value = global_dict.get(key, 0)
        if last_value == 0:
            global_dict.update({key: {'weight': value,
                                      'in_document': [index_doc]}})

        elif last_value['weight'] < value:

            global_dict[key]['weight'] = value
            global_dict[key]['in_document'].append(index_doc)

    def conveyor(self, remove_all_except_letter_dot, remove_stop_words):
        lemmatizer = WordNetLemmatizer()
        corpus = {}
        for name in os.listdir(self.directory):
            text = Extract(f"{self.directory}/{name}").extract()
            word_list = nltk.word_tokenize(text)

            text = ' '.join([lemmatizer.lemmatize(w) for w in word_list])

            if text is None:
                continue
            text = remove_stop_words(remove_all_except_letter_dot(text))
            corpus.update({name: text})

            print("=" * 20)
            self.result_extract_terms.update({name: self.get_terms_dict_with_text(text=text)})
        # pprint(corpus)
        self.corpus_dict = corpus

    @staticmethod
    def duplicate(global_dict, global_dict_keys):
        clean_terms = set()
        for el in global_dict_keys:
            tmp = el
            for el1 in global_dict_keys:
                if tmp == el1:
                    continue
                if len(set(tmp.lower().split()) & set(el1.lower().split())) > 0:
                    #                 print(tmp, len(tmp), 'vs', el1, len(el1))
                    if global_dict[el1]['weight'] > global_dict[tmp]['weight']:
                        tmp = el1
            clean_terms.add(tmp)
        if len(global_dict_keys) > len(clean_terms):
            return clean_terms, True
        else:
            return clean_terms, False

    @staticmethod
    def duplicate_lemms(list_terms, global_dict, lemma_dict):
        clean_terms = set()
        for k1 in list_terms:
            tmp = k1
            tmp_v1 = lemma_dict[k1]
            for k2 in list_terms:
                v2 = lemma_dict[k2]
                if tmp_v1 == v2:
                    continue
                if len(set(tmp_v1.lower().split()) & set(v2.lower().split())) > 0:
                    if global_dict[k2]['weight'] > global_dict[k1]['weight']:
                        tmp = k2
                        tmp_v1 = v2

            clean_terms.add(tmp)

        if len(list_terms) > len(clean_terms):
            return clean_terms, True
        else:
            return clean_terms, False

    @staticmethod
    def lemma_dict(nlp, text):
        lemma_dict = {}
        for sent in text:
            tmp = set()
            for word in nlp(sent.lower()):
                tmp.add(word.lemma_)
            lemma_dict.update({sent: " ".join(tmp)})
        return lemma_dict

    def cleaning(self):
        cleaning_dict = {}
        for index, dict_ in self.result_extract_terms.items():
            for key, value in dict_.items():
                self.update_score(global_dict=cleaning_dict, key=key, value=value, index_doc=index)
        self.cleaning_dict_terms = cleaning_dict


def remove_all_except_letter_dot_eng(text: str) -> str:
    return re.sub(r"[^A-Za-z.]", " ",
                  text)


def remove_stop_words_eng(text: str) -> str:
    stop_words = set(stopwords.words('english'))
    word_tokens = word_tokenize(text)
    filtered_sentence = []
    for w in word_tokens:
        if w not in stop_words:
            filtered_sentence.append(w)

    return " ".join(filtered_sentence)
